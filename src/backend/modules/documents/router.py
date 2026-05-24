import os
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from pydantic import BaseModel
from typing import List
from sqlalchemy.orm import Session
from src.backend.core.db import get_db
from src.backend.core.models import Document, Subject, User
from src.backend.modules.auth.dependencies import get_current_user

router = APIRouter(dependencies=[Depends(get_current_user)])

MAX_FILE_SIZE = 25 * 1024 * 1024
ALLOWED_EXTENSIONS = {".pdf", ".pptx"}


class DocumentOut(BaseModel):
    id: str
    filename: str
    status: str
    char_count: int | None = None
    created_at: str | None = None


class DocumentCreateOut(BaseModel):
    id: str
    filename: str
    status: str
    char_count: int | None = None


def _get_user_subject(subject_id: str, user: User, db: Session) -> Subject:
    subject = db.query(Subject).filter(
        Subject.id == subject_id,
        Subject.user_id == user.id,
    ).first()
    if not subject:
        raise HTTPException(status_code=404, detail="Not Found")
    return subject


def _get_extension(filename: str) -> str:
    _, ext = os.path.splitext(filename.lower())
    return ext


def _extract_text(filename: str, content: bytes) -> str:
    ext = _get_extension(filename)
    if ext == ".pdf":
        import fitz
        doc = fitz.open(stream=content, filetype="pdf")
        return "".join(page.get_text() for page in doc)
    elif ext == ".pptx":
        from io import BytesIO
        from pptx import Presentation
        prs = Presentation(BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "".join(texts)
    return ""


@router.get("/api/subjects/{subject_id}/documents", response_model=List[DocumentOut])
def list_documents(
    subject_id: str,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    subject = _get_user_subject(subject_id, user, db)
    docs = db.query(Document).filter(Document.subject_id == subject.id).all()
    return [
        DocumentOut(id=d.id, filename=d.filename, status=d.status,
                     char_count=d.char_count, created_at=d.created_at)
        for d in docs
    ]


@router.post("/api/subjects/{subject_id}/documents", status_code=201, response_model=List[DocumentCreateOut])
def create_document(
    subject_id: str,
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    subject = _get_user_subject(subject_id, user, db)
    results: List[DocumentCreateOut] = []

    for file in files:
        ext = _get_extension(file.filename or "")
        if ext not in ALLOWED_EXTENSIONS:
            results.append(DocumentCreateOut(id="", filename=file.filename or "unknown", status="error"))
            continue

        content = file.file.read()
        if len(content) > MAX_FILE_SIZE:
            results.append(DocumentCreateOut(id="", filename=file.filename or "unknown", status="error"))
            continue

        try:
            raw_text = _extract_text(file.filename or "", content)
        except Exception:
            results.append(DocumentCreateOut(id="", filename=file.filename or "unknown", status="error"))
            continue

        doc = Document(
            subject_id=subject.id,
            filename=file.filename or "unknown",
            raw_text=raw_text,
            char_count=len(raw_text),
            status="ready",
        )
        db.add(doc)
        db.commit()
        db.refresh(doc)
        results.append(DocumentCreateOut(id=doc.id, filename=doc.filename, status=doc.status, char_count=doc.char_count))

    return results


@router.delete("/api/documents/{document_id}", status_code=204)
def delete_document(
    document_id: str,
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
):
    doc = (
        db.query(Document)
        .join(Subject, Document.subject_id == Subject.id)
        .filter(Document.id == document_id, Subject.user_id == user.id)
        .first()
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Not Found")
    db.delete(doc)
    db.commit()