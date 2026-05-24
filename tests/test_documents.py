import pytest
from contextlib import contextmanager
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool
from src.backend.main import app
from src.backend.core.db import Base, get_db
from src.backend.core.models import User, Subject, Document
from src.backend.modules.auth.dependencies import get_current_user

engine = create_engine("sqlite://", connect_args={"check_same_thread": False}, poolclass=StaticPool)
TestingSessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)

fake_user = User(id="user-1", google_sub="sub-1", email="test@test.com", name="Test", created_at="2024-01-01")


def override_get_db():
    db = TestingSessionLocal()
    try:
        yield db
    finally:
        db.close()


def override_get_current_user():
    return fake_user


@contextmanager
def use_documents_overrides():
    app.dependency_overrides[get_db] = override_get_db
    app.dependency_overrides[get_current_user] = override_get_current_user
    try:
        yield
    finally:
        del app.dependency_overrides[get_db]
        del app.dependency_overrides[get_current_user]


def _make_minimal_pdf(text="Hello World from PDF"):
    content_stream_text = f"BT /F1 12 Tf 100 700 Td ({text}) Tj ET"
    content_stream = content_stream_text.encode("latin-1")
    stream_len = len(content_stream)

    parts = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (
            b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]\n"
            b"/Resources << /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >>\n"
            b"/Contents 4 0 R >>\nendobj\n"
        ),
        f"4 0 obj\n<< /Length {stream_len} >>\nstream\n{content_stream_text}\nendstream\nendobj\n".encode("latin-1"),
    ]

    header = b"%PDF-1.4\n"

    offsets = [0, len(header)]
    for i in range(len(parts) - 1):
        offsets.append(offsets[-1] + len(parts[i]))

    xref_entries = b"0000000000 65535 f \n"
    for i in range(1, len(parts) + 1):
        xref_entries += f"{offsets[i]:010d} 00000 n \n".encode()

    body = b"".join(parts)
    xref_offset = len(header) + len(body)
    xref = f"xref\n0 {len(parts) + 1}\n".encode() + xref_entries
    trailer = f"trailer\n<< /Size {len(parts) + 1} /Root 1 0 R >>\n".encode()
    startxref = f"startxref\n{xref_offset}\n%%EOF".encode()

    return header + body + xref + trailer + startxref


def _make_minimal_pptx():
    from pptx import Presentation
    from io import BytesIO
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.shapes.add_textbox(0, 0, 100, 50).text_frame.text = "Hello from PPTX"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestDocuments:

    def setup_method(self):
        Base.metadata.create_all(bind=engine)
        db = TestingSessionLocal()
        db.add(User(id="user-1", google_sub="sub-1", email="test@test.com", name="Test", created_at="2024-01-01"))
        db.add(Subject(id="subj-1", user_id="user-1", name="Math", created_at="2024-01-01"))
        db.commit()
        db.close()

    def teardown_method(self):
        Base.metadata.drop_all(bind=engine)

    def test_upload_pdf_returns_201_with_text_extracted(self, client):
        pdf_bytes = _make_minimal_pdf("Hello World from PDF")
        with use_documents_overrides():
            resp = client.post(
                "/api/subjects/subj-1/documents",
                files={"file": ("test.pdf", pdf_bytes, "application/pdf")},
                headers={"Authorization": "Bearer test"},
            )
        assert resp.status_code == 201, resp.text
        data = resp.json()
        assert data["status"] == "ready"
        assert data["char_count"] > 0
        assert data["filename"] == "test.pdf"
        assert "id" in data

    def test_upload_pptx_returns_201(self, client):
        pptx_bytes = _make_minimal_pptx()
        with use_documents_overrides():
            resp = client.post(
                "/api/subjects/subj-1/documents",
                files={"file": ("test.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                headers={"Authorization": "Bearer test"},
            )
        assert resp.status_code == 201, resp.text
        data = resp.json()
        assert data["status"] == "ready"
        assert data["char_count"] > 0

    def test_upload_txt_returns_400(self, client):
        with use_documents_overrides():
            resp = client.post(
                "/api/subjects/subj-1/documents",
                files={"file": ("notes.txt", b"some text", "text/plain")},
                headers={"Authorization": "Bearer test"},
            )
        assert resp.status_code == 400

    def test_upload_to_other_users_subject_returns_404(self, client):
        db = TestingSessionLocal()
        db.add(User(id="user-2", google_sub="sub-2", email="other@test.com", name="Other", created_at="2024-01-01"))
        db.add(Subject(id="subj-other", user_id="user-2", name="Physics", created_at="2024-01-01"))
        db.commit()
        db.close()
        pdf_bytes = _make_minimal_pdf("test")
        with use_documents_overrides():
            resp = client.post(
                "/api/subjects/subj-other/documents",
                files={"file": ("test.pdf", pdf_bytes, "application/pdf")},
                headers={"Authorization": "Bearer test"},
            )
        assert resp.status_code == 404

    def test_list_documents(self, client):
        db = TestingSessionLocal()
        db.add(Document(id="doc-1", subject_id="subj-1", filename="test.pdf",
                        status="ready", char_count=100, created_at="2024-01-01"))
        db.commit()
        db.close()
        with use_documents_overrides():
            resp = client.get(
                "/api/subjects/subj-1/documents",
                headers={"Authorization": "Bearer test"},
            )
        assert resp.status_code == 200
        data = resp.json()
        assert len(data) == 1
        assert data[0]["filename"] == "test.pdf"
        assert data[0]["status"] == "ready"

    def test_list_documents_scoped_to_subject(self, client):
        db = TestingSessionLocal()
        db.add(Document(id="doc-1", subject_id="subj-1", filename="doc1.pdf",
                        status="ready", char_count=100, created_at="2024-01-01"))
        db.add(Document(id="doc-2", subject_id="subj-1", filename="doc2.pdf",
                        status="ready", char_count=200, created_at="2024-01-01"))
        db.commit()
        db.close()
        with use_documents_overrides():
            resp = client.get(
                "/api/subjects/subj-1/documents",
                headers={"Authorization": "Bearer test"},
            )
        assert resp.status_code == 200
        assert len(resp.json()) == 2

    def test_delete_own_document(self, client):
        db = TestingSessionLocal()
        db.add(Document(id="doc-1", subject_id="subj-1", filename="test.pdf",
                        status="ready", char_count=100, created_at="2024-01-01"))
        db.commit()
        db.close()
        with use_documents_overrides():
            resp = client.delete("/api/documents/doc-1", headers={"Authorization": "Bearer test"})
        assert resp.status_code == 204
        db = TestingSessionLocal()
        assert db.query(Document).count() == 0
        db.close()

    def test_delete_nonexistent_document_returns_404(self, client):
        with use_documents_overrides():
            resp = client.delete("/api/documents/nonexistent", headers={"Authorization": "Bearer test"})
        assert resp.status_code == 404

    def test_delete_other_users_document_returns_404(self, client):
        db = TestingSessionLocal()
        db.add(User(id="user-2", google_sub="sub-2", email="other@test.com", name="Other", created_at="2024-01-01"))
        db.add(Subject(id="subj-other", user_id="user-2", name="Physics", created_at="2024-01-01"))
        db.add(Document(id="doc-other", subject_id="subj-other", filename="other.pdf",
                        status="ready", char_count=100, created_at="2024-01-01"))
        db.commit()
        db.close()
        with use_documents_overrides():
            resp = client.delete("/api/documents/doc-other", headers={"Authorization": "Bearer test"})
        assert resp.status_code == 404